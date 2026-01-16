import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import date, datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import re
import unicodedata
import requests
import base64
import os

# Document extraction libraries
try:
    import pdfplumber
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

# Translation library
try:
    from deep_translator import GoogleTranslator
    TRANSLATION_SUPPORT = True
except ImportError:
    TRANSLATION_SUPPORT = False

# Language codes mapping (used across translation functions)
LANG_CODES = {
    "Arabic": "ar", "English": "en", "French": "fr", "Spanish": "es",
    "German": "de", "Italian": "it", "Portuguese": "pt", "Russian": "ru",
    "Chinese (Simplified)": "zh-CN", "Chinese (Traditional)": "zh-TW",
    "Japanese": "ja", "Korean": "ko", "Hindi": "hi", "Turkish": "tr",
    "Dutch": "nl", "Polish": "pl", "Swedish": "sv", "Indonesian": "id",
    "Thai": "th", "Vietnamese": "vi", "Hebrew": "he", "Persian": "fa",
    "Urdu": "ur", "Bengali": "bn", "Greek": "el"
}

# RTL (Right-to-Left) languages
RTL_LANGUAGES = {"Arabic", "Hebrew", "Persian", "Urdu"}

# Free Cloud LLM Support (HuggingFace free tier - no API key needed)
LLM_AVAILABLE = False

# Mistral-7B-Instruct - Better for large texts and summaries (8K context)
LLM_MODEL = "mistralai/Mistral-7B-Instruct-v0.2"

# Use HuggingFace Inference API (free tier, no key required)
try:
    test_response = requests.post(
        f"https://api-inference.huggingface.co/models/{LLM_MODEL}",
        json={"inputs": "test"},
        timeout=5
    )
    if test_response.status_code in [200, 503]:  # 503 means model loading
        LLM_AVAILABLE = True
except Exception:
    pass

# --- TEXT CLEANING FUNCTIONS ---
def clean_extracted_text(text):
    """Clean and normalize extracted text, handling mixed language concatenation"""
    if not text:
        return ""
    
    # Normalize unicode characters
    text = unicodedata.normalize('NFKC', text)
    
    # Remove zero-width characters first
    text = re.sub(r'[\u200b\u200c\u200d\ufeff\u00ad]', '', text)
    
    # Remove control characters except newlines and tabs
    text = ''.join(char for char in text if unicodedata.category(char) != 'Cc' or char in '\n\t')
    
    # Add space between Arabic and Latin characters (handles concatenation)
    text = re.sub(r'([\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF])([A-Za-z0-9])', r'\1 | \2', text)
    text = re.sub(r'([A-Za-z0-9])([\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF])', r'\1 | \2', text)
    
    # Standardize quotes
    text = text.replace('"', '"').replace('"', '"').replace(''', "'").replace(''', "'")
    
    # Clean up multiple spaces
    text = re.sub(r'[ \t]+', ' ', text)
    
    # Clean up multiple newlines (keep max 2)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Clean up lines
    lines = text.split('\n')
    cleaned_lines = []
    for line in lines:
        line = line.strip()
        if line:
            cleaned_lines.append(line)
    
    return '\n'.join(cleaned_lines)

def extract_english_only(text):
    """Extract only English text from mixed content"""
    parts = re.split(r'\s*\|\s*', text)
    english_parts = []
    for part in parts:
        if part.strip() and not re.search(r'[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF]', part):
            english_parts.append(part.strip())
    return ' '.join(english_parts)

# --- LLM HELPER FUNCTIONS ---
def query_llm(prompt, max_tokens=1024):
    """Query HuggingFace free tier LLM (no API key needed) - Mistral-7B for better summaries"""
    if not LLM_AVAILABLE:
        return None
    
    try:
        # Format prompt for Mistral instruction format
        formatted_prompt = f"<s>[INST] {prompt} [/INST]"
        
        # Use free HuggingFace inference API with Mistral-7B-Instruct
        response = requests.post(
            f"https://api-inference.huggingface.co/models/{LLM_MODEL}",
            json={
                "inputs": formatted_prompt, 
                "parameters": {
                    "max_new_tokens": max_tokens,
                    "temperature": 0.3,
                    "do_sample": True,
                    "return_full_text": False
                }
            },
            timeout=120
        )
        if response.status_code == 200:
            result = response.json()
            if isinstance(result, list) and len(result) > 0:
                return result[0].get('generated_text', '')
            elif isinstance(result, dict) and 'generated_text' in result:
                return result['generated_text']
            return str(result)
        elif response.status_code == 503:
            # Model is loading, wait and retry once
            import time
            time.sleep(30)
            response = requests.post(
                f"https://api-inference.huggingface.co/models/{LLM_MODEL}",
                json={
                    "inputs": formatted_prompt, 
                    "parameters": {
                        "max_new_tokens": max_tokens,
                        "temperature": 0.3,
                        "do_sample": True,
                        "return_full_text": False
                    }
                },
                timeout=120
            )
            if response.status_code == 200:
                result = response.json()
                if isinstance(result, list) and len(result) > 0:
                    return result[0].get('generated_text', '')
        return None
    except Exception:
        return None

# --- DOCUMENT EXTRACTION FUNCTIONS ---
def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return clean_extracted_text(text)

def extract_text_from_docx(file):
    """Extract text from DOCX file"""
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return clean_extracted_text(text)

def extract_text_from_pptx(file):
    """Extract text from PPTX file"""
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_extracted_text(text)

def extract_text(uploaded_file):
    """Extract text based on file type"""
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    if file_type == 'pdf':
        if not PDF_SUPPORT:
            return None, "PDF support not installed. Run: pip install pdfplumber"
        return extract_text_from_pdf(uploaded_file), None
    elif file_type == 'docx':
        if not DOCX_SUPPORT:
            return None, "DOCX support not installed. Run: pip install python-docx"
        return extract_text_from_docx(uploaded_file), None
    elif file_type in ['pptx', 'ppt']:
        return extract_text_from_pptx(uploaded_file), None
    else:
        return None, f"Unsupported file type: {file_type}"

# --- TRANSLATION FUNCTION ---
def translate_text(text, target_lang):
    """Translate text to target language"""
    if not TRANSLATION_SUPPORT:
        return "Translation not available. Install: pip install deep-translator"
    
    target_code = LANG_CODES.get(target_lang, "en")
    
    # Clean text before translation
    text = clean_extracted_text(text)
    
    # Split text into chunks (Google Translate has a 5000 char limit)
    max_chars = 4500
    chunks = [text[i:i+max_chars] for i in range(0, len(text), max_chars)]
    
    translated_chunks = []
    for chunk in chunks:
        if chunk.strip():
            try:
                translated = GoogleTranslator(source='auto', target=target_code).translate(chunk)
                translated_chunks.append(translated if translated else chunk)
            except Exception as e:
                translated_chunks.append(f"[Translation error: {str(e)}]")
    
    return "\n".join(translated_chunks)

def translate_single_text(text, target_code):
    """Translate a single text chunk"""
    if not text or not text.strip():
        return text
    if not TRANSLATION_SUPPORT:
        return text
    try:
        # Remove pipe separators before translation
        clean_text = text.replace(' | ', ' ').strip()
        if len(clean_text) > 4500:
            clean_text = clean_text[:4500]
        translated = GoogleTranslator(source='auto', target=target_code).translate(clean_text)
        return translated if translated else text
    except Exception:
        return text

def translate_docx_inplace(file, target_lang):
    """Translate a DOCX file preserving original structure and formatting"""
    if not DOCX_SUPPORT:
        return None, "DOCX support not installed"
    
    target_code = LANG_CODES.get(target_lang, "en")
    
    # Load the document
    doc = Document(file)
    
    # Translate paragraphs while preserving formatting
    for para in doc.paragraphs:
        if para.text.strip():
            # Translate full paragraph text
            translated = translate_single_text(para.text, target_code)
            
            # Clear and set new text (preserves paragraph-level formatting)
            if para.runs:
                # Keep first run's formatting, update text
                first_run = para.runs[0]
                original_font = first_run.font.name
                original_size = first_run.font.size
                original_bold = first_run.font.bold
                original_italic = first_run.font.italic
                
                # Clear all runs
                for run in para.runs:
                    run.text = ""
                
                # Set translated text in first run
                first_run.text = translated
                first_run.font.name = original_font
                first_run.font.size = original_size
                first_run.font.bold = original_bold
                first_run.font.italic = original_italic
            else:
                para.text = translated
            
            # Set RTL alignment for RTL languages
            if target_lang in RTL_LANGUAGES:
                para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    
    # Translate tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if para.text.strip():
                        translated = translate_single_text(para.text, target_code)
                        if para.runs:
                            para.runs[0].text = translated
                            for run in para.runs[1:]:
                                run.text = ""
                        else:
                            para.text = translated
    
    # Save to buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer, None

def translate_pptx_inplace(file, target_lang):
    """Translate a PPTX file preserving original structure and formatting"""
    target_code = LANG_CODES.get(target_lang, "en")
    
    # Load presentation
    prs = Presentation(file)
    
    # Translate each slide
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            # Preserve formatting
                            original_font = run.font.name
                            original_size = run.font.size
                            original_bold = run.font.bold
                            original_italic = run.font.italic
                            
                            # Safely get color - handle scheme colors vs RGB
                            original_color = None
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    try:
                                        original_color = run.font.color.rgb
                                    except AttributeError:
                                        # It's a scheme/theme color, skip RGB
                                        original_color = None
                            except Exception:
                                original_color = None
                            
                            # Translate
                            run.text = translate_single_text(run.text, target_code)
                            
                            # Restore formatting
                            run.font.name = original_font
                            run.font.size = original_size
                            run.font.bold = original_bold
                            run.font.italic = original_italic
                            if original_color:
                                run.font.color.rgb = original_color
            
            # Handle tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.text.strip():
                                        run.text = translate_single_text(run.text, target_code)
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer, None

def create_translated_docx(translated_text, target_lang, original_filename):
    """Create a DOCX document from translated text"""
    doc = Document()
    
    # Add title
    doc.add_heading(f'Translated Document ({target_lang})', 0)
    
    # Add metadata
    doc.add_paragraph(f"Original file: {original_filename}")
    doc.add_paragraph(f"Translation date: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    doc.add_paragraph("---")
    
    # Add translated content
    paragraphs = translated_text.split('\n')
    for para in paragraphs:
        if para.strip():
            p = doc.add_paragraph(para.strip())
            # Set RTL for RTL languages
            if target_lang in RTL_LANGUAGES:
                p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    
    # Save to bytes
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_translated_pptx(translated_text, target_lang, original_filename):
    """Create a PPTX presentation from translated text"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Translated Document ({target_lang})"
    subtitle.text = f"Original: {original_filename}\nDate: {datetime.now().strftime('%Y-%m-%d')}"
    
    # Content slides
    paragraphs = [p.strip() for p in translated_text.split('\n') if p.strip()]
    content_layout = prs.slide_layouts[1]
    
    # Group paragraphs into slides (max 6 paragraphs per slide)
    for i in range(0, len(paragraphs), 6):
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = f"Content (Page {i//6 + 1})"
        
        body = slide.placeholders[1]
        tf = body.text_frame
        
        for j, para in enumerate(paragraphs[i:i+6]):
            if j == 0:
                tf.text = para[:200]  # Limit text length
            else:
                p = tf.add_paragraph()
                p.text = para[:200]
    
    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# --- ANALYSIS FUNCTIONS ---
def generate_summary(text):
    """Generate a comprehensive, detailed summary with explanations from extracted text"""
    # Clean and prepare text - extract English for better summarization
    cleaned_text = clean_extracted_text(text)
    english_text = extract_english_only(cleaned_text)
    
    # If mostly Arabic, use the full cleaned text
    if len(english_text) < 100:
        work_text = cleaned_text.replace(' | ', ' ')
    else:
        work_text = english_text
    
    # Try LLM first for better summary
    if LLM_AVAILABLE and len(work_text) > 50:
        prompt = f"""Analyze this document in detail and create a comprehensive executive summary with explanations.

Structure your response as follows:

1. **Overview**: Provide a 2-3 sentence overview of what this document is about.

2. **Key Objectives**: List and explain the main objectives, goals, or purposes mentioned in the document. Provide context and reasoning for each.

3. **Main Content**: Summarize the core content in detail. Include important details, initiatives, strategies, or processes mentioned.

4. **Stakeholders & Resources**: Identify any mentioned teams, departments, stakeholders, or resources involved.

5. **Important Details**: Highlight any critical information, requirements, constraints, or special considerations.

6. **Implications**: Explain what this means for the organization or project and why it matters.

Document Content:
{work_text[:6000]}

Comprehensive Summary:"""
        llm_summary = query_llm(prompt, max_tokens=1024)
        if llm_summary and len(llm_summary) > 50:
            return llm_summary.strip()
    
    # Fallback to intelligent extraction (more comprehensive)
    # Look for key sections and extract meaningful content
    summary_parts = []
    
    # Extract sentences, handling multiple delimiters
    sentences = re.split(r'[.!?؟。\n]+', work_text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 15]
    
    # Prioritize sentences with key terms (expanded)
    priority_terms = ['objective', 'goal', 'strategy', 'initiative', 'platform', 
                      'department', 'ai', 'digital', 'transformation', 'key', 'main',
                      'purpose', 'vision', 'mission', 'target', 'deliver', 'implement',
                      'develop', 'create', 'establish', 'enhance', 'improve', 'focus']
    
    context_terms = ['background', 'overview', 'introduction', 'summary', 'scope',
                     'stakeholder', 'team', 'resource', 'requirement', 'constraint']
    
    detail_terms = ['include', 'feature', 'component', 'process', 'phase', 'stage',
                    'milestone', 'deliverable', 'output', 'result', 'outcome']
    
    priority_sentences = []
    context_sentences = []
    detail_sentences = []
    other_sentences = []
    
    for s in sentences:
        s_lower = s.lower()
        if any(term in s_lower for term in priority_terms):
            priority_sentences.append(s)
        elif any(term in s_lower for term in context_terms):
            context_sentences.append(s)
        elif any(term in s_lower for term in detail_terms):
            detail_sentences.append(s)
        else:
            other_sentences.append(s)
    
    # Build comprehensive summary with structure
    summary = "**📋 Document Overview:**\n"
    
    # Add overview sentences
    overview = context_sentences[:2] if context_sentences else other_sentences[:2]
    for sent in overview:
        sent = sent.strip()
        if not sent.endswith(('.', '!', '?')):
            sent += '.'
        summary += f"{sent}\n\n"
    
    # Add key objectives
    if priority_sentences:
        summary += "**🎯 Key Objectives & Focus Areas:**\n"
        for i, sent in enumerate(priority_sentences[:5], 1):
            sent = sent.strip()
            if not sent.endswith(('.', '!', '?')):
                sent += '.'
            summary += f"{i}. {sent}\n"
        summary += "\n"
    
    # Add important details
    if detail_sentences:
        summary += "**📌 Important Details:**\n"
        for sent in detail_sentences[:4]:
            sent = sent.strip()
            if not sent.endswith(('.', '!', '?')):
                sent += '.'
            summary += f"• {sent}\n"
        summary += "\n"
    
    # Add additional context
    if len(other_sentences) > 2:
        summary += "**ℹ️ Additional Context:**\n"
        for sent in other_sentences[2:5]:
            sent = sent.strip()
            if not sent.endswith(('.', '!', '?')):
                sent += '.'
            summary += f"• {sent}\n"
    
    if len(summary) > 50:
        return summary
    
    return "No comprehensive summary could be generated from the document content."

def extract_dates(text):
    """Extract dates mentioned in the document"""
    # Pattern for various date formats
    date_patterns = [
        r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}',  # DD/MM/YYYY or MM-DD-YYYY
        r'\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}',
        r'(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}',
    ]
    
    dates_found = []
    for pattern in date_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        dates_found.extend(matches)
    
    return list(set(dates_found))

def extract_keywords(text):
    """Extract important keywords/phrases using LLM if available"""
    # Try LLM first
    if LLM_AVAILABLE:
        prompt = f"""Extract the 10 most important keywords or key phrases from this document. Return them as a comma-separated list:

{text[:2000]}

Keywords:"""
        llm_keywords = query_llm(prompt, max_tokens=100)
        if llm_keywords:
            # Parse the response
            keywords = [k.strip() for k in llm_keywords.split(',') if k.strip()]
            if keywords:
                return keywords[:10]
    
    # Fallback: Common project-related keywords
    keywords = []
    important_terms = [
        'budget', 'timeline', 'deadline', 'milestone', 'phase', 'objective',
        'risk', 'requirement', 'stakeholder', 'deliverable', 'scope',
        'implementation', 'deployment', 'integration', 'compliance', 'governance',
        'strategy', 'digital', 'transformation', 'automation', 'efficiency',
        'innovation', 'performance', 'quality', 'security', 'data'
    ]
    
    text_lower = text.lower()
    for term in important_terms:
        if term in text_lower:
            keywords.append(term.capitalize())
    
    return keywords

def calculate_risk_score(text):
    """Calculate risk score based on document content using LLM if available"""
    # Try LLM first for better analysis
    if LLM_AVAILABLE:
        prompt = f"""Analyze this document for project risks. Rate the overall risk level as LOW, MEDIUM, or HIGH, and provide a score from 1-10 (10 being lowest risk). Format: "LEVEL: X/10"

{text[:2000]}

Risk Assessment:"""
        llm_response = query_llm(prompt, max_tokens=100)
        if llm_response:
            # Parse response
            response_lower = llm_response.lower()
            if 'high' in response_lower:
                level = "High"
                score = 3.0
            elif 'low' in response_lower:
                level = "Low"
                score = 8.0
            else:
                level = "Medium"
                score = 5.0
            
            # Try to extract numeric score
            score_match = re.search(r'(\d+(?:\.\d+)?)\s*/\s*10', llm_response)
            if score_match:
                score = float(score_match.group(1))
            
            return round(score, 1), level
    
    # Fallback to rule-based
    risk_indicators = {
        'high_risk': ['urgent', 'critical', 'risk', 'delay', 'issue', 'problem', 'challenge', 'concern', 'failure', 'crisis'],
        'medium_risk': ['consider', 'review', 'assess', 'evaluate', 'potential', 'uncertain', 'unclear'],
        'low_risk': ['complete', 'success', 'achieved', 'approved', 'confirmed', 'stable', 'secure']
    }
    
    text_lower = text.lower()
    high_count = sum(1 for word in risk_indicators['high_risk'] if word in text_lower)
    medium_count = sum(1 for word in risk_indicators['medium_risk'] if word in text_lower)
    low_count = sum(1 for word in risk_indicators['low_risk'] if word in text_lower)
    
    # Calculate weighted score (lower is better)
    total = high_count * 3 + medium_count * 2 + low_count * 1
    if total == 0:
        return 5.0, "Medium"
    
    score = min(10, max(1, 10 - (high_count * 2) + (low_count * 0.5)))
    
    if score >= 7:
        level = "Low"
    elif score >= 4:
        level = "Medium"
    else:
        level = "High"
    
    return round(score, 1), level

def estimate_project_duration(text):
    """Estimate project duration based on content"""
    # Look for duration mentions
    duration_patterns = [
        (r'(\d+)\s*weeks?', 'weeks'),
        (r'(\d+)\s*months?', 'months'),
        (r'(\d+)\s*days?', 'days'),
    ]
    
    for pattern, unit in duration_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            num = int(match.group(1))
            return f"{num} {unit.capitalize()}"
    
    # Default estimate based on document length
    word_count = len(text.split())
    if word_count > 2000:
        return "12-16 Weeks"
    elif word_count > 1000:
        return "8-12 Weeks"
    else:
        return "4-8 Weeks"

def generate_timeline(text):
    """Generate project timeline from document using AI when available"""
    today = datetime.now()
    timeline_data = []
    
    # Try LLM first for intelligent phase extraction
    if LLM_AVAILABLE:
        prompt = f"""Analyze this project document and extract the project phases/stages with estimated durations.
Format each phase as: "Phase Name | Duration in days"
List 4-6 phases. If durations aren't mentioned, estimate based on complexity.

Document:
{text[:2500]}

Project Phases:"""
        
        llm_response = query_llm(prompt, max_tokens=300)
        
        if llm_response:
            # Parse LLM response
            lines = llm_response.strip().split('\n')
            current_date = today
            
            for line in lines:
                line = line.strip()
                if not line or line.startswith('#'):
                    continue
                
                # Remove bullet points, numbers at start
                line = re.sub(r'^[\d\.\-\*\•]+\s*', '', line)
                
                # Try to parse "Phase Name | Duration" format
                if '|' in line:
                    parts = line.split('|')
                    phase_name = parts[0].strip()
                    duration_text = parts[1].strip() if len(parts) > 1 else "14"
                    
                    # Extract number from duration
                    duration_match = re.search(r'(\d+)', duration_text)
                    duration = int(duration_match.group(1)) if duration_match else 14
                    
                    # Cap duration to reasonable range
                    duration = max(7, min(90, duration))
                else:
                    # Just phase name, use default duration
                    phase_name = line
                    duration = 14
                
                if phase_name and len(phase_name) > 3:
                    # Determine resource based on phase content
                    phase_lower = phase_name.lower()
                    if any(word in phase_lower for word in ['plan', 'requirement', 'analysis', 'design', 'scope']):
                        resource = "Planning"
                    elif any(word in phase_lower for word in ['test', 'valid', 'qa', 'quality', 'review']):
                        resource = "QA"
                    elif any(word in phase_lower for word in ['deploy', 'handover', 'launch', 'release', 'go-live']):
                        resource = "Operations"
                    elif any(word in phase_lower for word in ['develop', 'implement', 'build', 'code', 'create']):
                        resource = "Development"
                    else:
                        resource = "Execution"
                    
                    start = current_date
                    finish = current_date + timedelta(days=duration)
                    
                    timeline_data.append({
                        'Task': phase_name[:60],  # Limit length
                        'Start': start.strftime('%Y-%m-%d'),
                        'Finish': finish.strftime('%Y-%m-%d'),
                        'Resource': resource
                    })
                    
                    current_date = finish + timedelta(days=1)
            
            # If we got valid phases from LLM, return them
            if len(timeline_data) >= 2:
                return pd.DataFrame(timeline_data)
    
    # Fallback: Try to find phases mentioned in text using regex
    phases = []
    phase_patterns = [
        r'phase\s*(\d+)[:\s]*([^\n.]+)',
        r'step\s*(\d+)[:\s]*([^\n.]+)',
        r'stage\s*(\d+)[:\s]*([^\n.]+)',
        r'milestone\s*(\d+)[:\s]*([^\n.]+)',
    ]
    
    for pattern in phase_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        for match in matches:
            phases.append(f"Phase {match[0]}: {match[1].strip()[:50]}")
    
    # Also look for numbered lists that might be phases
    numbered_items = re.findall(r'^\s*(\d+)[\.\)]\s*([A-Z][^\n]{10,60})', text, re.MULTILINE)
    if len(numbered_items) >= 3 and not phases:
        for num, item in numbered_items[:6]:
            phases.append(f"Step {num}: {item.strip()}")
    
    # If still no phases found, create default structure based on document content
    if not phases:
        # Analyze content to create relevant default phases
        text_lower = text.lower()
        
        if 'digital' in text_lower or 'transformation' in text_lower:
            phases = [
                "Phase 1: Assessment & Strategy",
                "Phase 2: Digital Infrastructure Setup",
                "Phase 3: Implementation & Integration",
                "Phase 4: Testing & Optimization",
                "Phase 5: Deployment & Training"
            ]
        elif 'software' in text_lower or 'development' in text_lower:
            phases = [
                "Phase 1: Requirements Gathering",
                "Phase 2: Design & Architecture",
                "Phase 3: Development",
                "Phase 4: Testing & QA",
                "Phase 5: Deployment & Support"
            ]
        else:
            phases = [
                "Phase 1: Initiation & Planning",
                "Phase 2: Analysis & Design",
                "Phase 3: Execution",
                "Phase 4: Monitoring & Control",
                "Phase 5: Closure & Handover"
            ]
    
    # Generate timeline data with varied durations
    current_date = today
    base_duration = 14  # 2 weeks base
    
    for i, phase in enumerate(phases[:6]):  # Max 6 phases
        # Vary duration based on phase type
        phase_lower = phase.lower()
        if 'plan' in phase_lower or 'requirement' in phase_lower or 'analysis' in phase_lower:
            duration = base_duration  # 2 weeks for planning
        elif 'develop' in phase_lower or 'implement' in phase_lower or 'execution' in phase_lower:
            duration = base_duration * 2  # 4 weeks for development
        elif 'test' in phase_lower or 'qa' in phase_lower:
            duration = int(base_duration * 1.5)  # 3 weeks for testing
        else:
            duration = base_duration
        
        start = current_date
        finish = current_date + timedelta(days=duration)
        
        # Determine resource based on phase
        if 'plan' in phase_lower or 'requirement' in phase_lower or 'analysis' in phase_lower:
            resource = "Planning"
        elif 'test' in phase_lower or 'valid' in phase_lower or 'qa' in phase_lower:
            resource = "QA"
        elif 'deploy' in phase_lower or 'handover' in phase_lower or 'closure' in phase_lower:
            resource = "Operations"
        elif 'develop' in phase_lower or 'implement' in phase_lower:
            resource = "Development"
        else:
            resource = "Execution"
        
        timeline_data.append({
            'Task': phase,
            'Start': start.strftime('%Y-%m-%d'),
            'Finish': finish.strftime('%Y-%m-%d'),
            'Resource': resource
        })
        
        current_date = finish + timedelta(days=1)
    
    return pd.DataFrame(timeline_data)

def analyze_go_nogo(text):
    """Perform Go/No-Go analysis using LLM for intelligent decision making"""
    criteria = {}
    ai_reasoning = ""
    ai_verdict = None
    
    # Try LLM for comprehensive AI-driven analysis
    if LLM_AVAILABLE:
        # First prompt: Get detailed scores with reasoning
        score_prompt = f"""Analyze this project document for a Go/No-Go decision.

Rate each criterion from 1-10 (1=very poor, 10=excellent):
- Technical Feasibility: Can the project be technically implemented?
- Budget Availability: Are financial resources adequate?
- Resource Readiness: Are team and skills available?
- Stakeholder Alignment: Is there executive support?

Format your response as:
Technical Feasibility: X/10
Budget Availability: X/10
Resource Readiness: X/10
Stakeholder Alignment: X/10
Verdict: GO or NO-GO or CONDITIONAL

Document:
{text[:2500]}

Analysis:"""
        
        llm_response = query_llm(score_prompt, max_tokens=300)
        
        if llm_response:
            ai_reasoning = llm_response
            
            # Parse scores using multiple patterns for robustness
            score_patterns = [
                (r'Technical\s*Feasibility[:\s]*(\d+)', 'Technical Feasibility'),
                (r'Budget\s*Availability[:\s]*(\d+)', 'Budget Availability'),
                (r'Resource\s*Readiness[:\s]*(\d+)', 'Resource Readiness'),
                (r'Stakeholder\s*Alignment[:\s]*(\d+)', 'Stakeholder Alignment'),
            ]
            
            for pattern, criterion in score_patterns:
                match = re.search(pattern, llm_response, re.IGNORECASE)
                if match:
                    score = int(match.group(1))
                    criteria[criterion] = min(10, max(1, score))
            
            # Also try simpler number extraction if structured format fails
            if len(criteria) < 4:
                # Look for any numbers in the response
                numbers = re.findall(r'\b(\d+)\s*/\s*10', llm_response)
                if not numbers:
                    numbers = re.findall(r'\b([1-9]|10)\b', llm_response)
                
                criterion_names = ['Technical Feasibility', 'Budget Availability', 'Resource Readiness', 'Stakeholder Alignment']
                for i, criterion in enumerate(criterion_names):
                    if criterion not in criteria and i < len(numbers):
                        criteria[criterion] = min(10, max(1, int(numbers[i])))
            
            # Parse verdict from AI response
            response_lower = llm_response.lower()
            if 'no-go' in response_lower or 'nogo' in response_lower or 'not recommended' in response_lower:
                ai_verdict = "NO-GO"
            elif 'conditional' in response_lower or 'with conditions' in response_lower or 'pending' in response_lower:
                ai_verdict = "CONDITIONAL"
            elif 'go' in response_lower and 'no-go' not in response_lower:
                ai_verdict = "GO"
    
    # Fill in missing criteria with AI-informed fallback
    text_lower = text.lower()
    
    if 'Technical Feasibility' not in criteria:
        # Analyze technical indicators
        tech_positive = ['proven', 'established', 'available', 'existing', 'ready', 'mature', 'stable']
        tech_negative = ['complex', 'new technology', 'untested', 'experimental', 'challenging', 'difficult']
        positive_count = sum(1 for k in tech_positive if k in text_lower)
        negative_count = sum(1 for k in tech_negative if k in text_lower)
        tech_keywords = ['technical', 'technology', 'system', 'platform', 'infrastructure', 'integration', 'ai', 'digital']
        base_score = 5 + min(3, sum(1 for k in tech_keywords if k in text_lower))
        criteria['Technical Feasibility'] = min(10, max(1, base_score + positive_count - negative_count))
    
    if 'Budget Availability' not in criteria:
        # Analyze budget indicators
        budget_positive = ['funded', 'approved', 'allocated', 'sufficient', 'available', 'secured']
        budget_negative = ['limited', 'constraint', 'shortage', 'insufficient', 'unfunded', 'pending approval']
        positive_count = sum(1 for k in budget_positive if k in text_lower)
        negative_count = sum(1 for k in budget_negative if k in text_lower)
        budget_keywords = ['budget', 'cost', 'fund', 'invest', 'financial', 'capital']
        base_score = 5 + min(2, sum(1 for k in budget_keywords if k in text_lower))
        criteria['Budget Availability'] = min(10, max(1, base_score + positive_count - negative_count))
    
    if 'Resource Readiness' not in criteria:
        # Analyze resource indicators
        resource_positive = ['experienced', 'skilled', 'trained', 'qualified', 'available', 'dedicated', 'capable']
        resource_negative = ['shortage', 'hiring', 'training needed', 'gap', 'lack', 'insufficient']
        positive_count = sum(1 for k in resource_positive if k in text_lower)
        negative_count = sum(1 for k in resource_negative if k in text_lower)
        resource_keywords = ['team', 'staff', 'resource', 'personnel', 'expert', 'skill']
        base_score = 4 + min(2, sum(1 for k in resource_keywords if k in text_lower))
        criteria['Resource Readiness'] = min(10, max(1, base_score + positive_count - negative_count))
    
    if 'Stakeholder Alignment' not in criteria:
        # Analyze stakeholder indicators
        stakeholder_positive = ['approved', 'supported', 'endorsed', 'committed', 'aligned', 'agreed']
        stakeholder_negative = ['opposition', 'concern', 'resistance', 'disagreement', 'pending', 'unclear']
        positive_count = sum(1 for k in stakeholder_positive if k in text_lower)
        negative_count = sum(1 for k in stakeholder_negative if k in text_lower)
        stakeholder_keywords = ['stakeholder', 'sponsor', 'executive', 'management', 'leadership', 'board']
        base_score = 5 + min(2, sum(1 for k in stakeholder_keywords if k in text_lower))
        criteria['Stakeholder Alignment'] = min(10, max(1, base_score + positive_count - negative_count))
    
    # Calculate overall score and determine verdict (forgiving thresholds)
    avg_score = sum(criteria.values()) / len(criteria)
    
    # Use AI verdict if available, otherwise calculate based on scores
    if ai_verdict:
        verdict = f"{ai_verdict} {'✅' if ai_verdict == 'GO' else '❌' if ai_verdict == 'NO-GO' else '⚠️'}"
        # Adjust confidence based on score alignment with AI verdict
        if ai_verdict == "GO" and avg_score >= 5:
            confidence = int(min(95, avg_score * 10 + 15))
        elif ai_verdict == "NO-GO" and avg_score < 3:
            confidence = int(min(95, (10 - avg_score) * 10 + 10))
        elif ai_verdict == "CONDITIONAL":
            confidence = int(avg_score * 10 + 10)
        else:
            confidence = int(avg_score * 9)  # More forgiving when AI and scores disagree
    else:
        # Score-based verdict (more forgiving thresholds)
        if avg_score >= 5:  # Lowered from 7 - easier to get GO
            verdict = "GO ✅"
            confidence = int(min(95, avg_score * 10 + 10))
        elif avg_score >= 3:  # Lowered from 5 - wider CONDITIONAL range
            verdict = "CONDITIONAL ⚠️"
            confidence = int(avg_score * 10 + 15)
        else:  # Only NO-GO if really low (below 3)
            verdict = "NO-GO ❌"
            confidence = int((10 - avg_score) * 7)
    
    return criteria, verdict, confidence

# --- 1. PAGE CONFIGURATION ---
st.set_page_config(
    page_title="Digital Transformation Hub",
    page_icon="🟩",
    layout="wide"
)

# --- 2. LOAD CSS AND PATTERN BORDERS ---
def load_css():
    with open("style.css") as f:
        css_content = f.read()
    
    st.markdown(f"<style>{css_content}</style>", unsafe_allow_html=True)

def load_pattern_borders():
    """Inject pattern border elements as HTML divs"""
    pattern_path = "pattern.png"
    if os.path.exists(pattern_path):
        with open(pattern_path, "rb") as img_file:
            pattern_base64 = base64.b64encode(img_file.read()).decode()
        
        pattern_html = f'''
        <div class="pattern-border pattern-border-left" 
             style="background-image: url('data:image/png;base64,{pattern_base64}');"></div>
        <div class="pattern-border pattern-border-right" 
             style="background-image: url('data:image/png;base64,{pattern_base64}');"></div>
        <div class="pattern-border-horizontal pattern-border-top" 
             style="background-image: url('data:image/png;base64,{pattern_base64}');"></div>
        <div class="pattern-border-horizontal pattern-border-bottom" 
             style="background-image: url('data:image/png;base64,{pattern_base64}');"></div>
        '''
        st.markdown(pattern_html, unsafe_allow_html=True)

load_css()
load_pattern_borders()

# --- 3. MAIN INTERFACE ---

# Header Section with logos (always visible)
col1, col2 = st.columns([5, 1])
with col1:
    st.image("https://placehold.co/250x80/000000/FFFFFF?text=Intelligent+Solution", width=250)
    st.markdown("**Digital Transformation Enabler** | Fast-Track Solution Rebuild")
with col2:
    # DC Logo placeholder on the right
    st.image("icon.png", width=100)

# Initialize session state for user ID
if 'user_id' not in st.session_state:
    st.session_state.user_id = None

# User ID Entry Screen
if st.session_state.user_id is None:
    st.markdown("---")
    
    # Center the ID input using columns
    col_left, col_center, col_right = st.columns([1, 1, 1])
    with col_center:
        st.markdown("### Please enter your ID")
        with st.form("user_id_form", clear_on_submit=False):
            st.markdown('<div class="id-input-container">', unsafe_allow_html=True)
            user_id_input = st.text_input("Enter your ID", placeholder="Enter any ID...", label_visibility="collapsed")
            st.markdown('</div>', unsafe_allow_html=True)
            submitted = st.form_submit_button("Continue", type="primary", use_container_width=True)
            if submitted:
                if user_id_input.strip():
                    st.session_state.user_id = user_id_input.strip()
                    st.rerun()
                else:
                    st.error("Please enter a valid ID")
    
    st.stop()  # Stop execution here until ID is entered

# Language options - derive from LANG_CODES constant
languages = list(LANG_CODES.keys())

# Main Content Area
st.markdown("### 1. Document Ingestion")

col_left, col_upload, col_right = st.columns([1, 2, 1])
with col_upload:
    uploaded_file = st.file_uploader("Upload Document", type=['pdf', 'docx', 'pptx', 'ppt'], label_visibility="collapsed")
    st.markdown("**Target language**")
    target_lang = st.selectbox("Target Language", languages, label_visibility="collapsed")

if uploaded_file:
    # REAL DOCUMENT PROCESSING
    with st.spinner('Extracting text from document...'):
        extracted_text, error = extract_text(uploaded_file)
    
    if error:
        st.error(error)
    elif not extracted_text or len(extracted_text.strip()) < 10:
        st.warning("Could not extract meaningful text from the document. Please check the file.")
    else:
        # Text is already cleaned by extract functions, but ensure it's clean
        extracted_text = clean_extracted_text(extracted_text)
        
        
        st.markdown("---")
        
        # --- REAL ANALYSIS ---
        with st.spinner('Analyzing document content...'):
            # Calculate all metrics
            risk_score, risk_level = calculate_risk_score(extracted_text)
            duration = estimate_project_duration(extracted_text)
            keywords = extract_keywords(extracted_text)
            dates_found = extract_dates(extracted_text)
            criteria, verdict, confidence = analyze_go_nogo(extracted_text)
        
        # --- DASHBOARD LAYOUT ---
        st.markdown("### 2. Strategic Analysis")
        
        # Main Navigation Tabs (replacing metric boxes)
        tab_summary, tab_timeline, tab_decision, tab_translation, tab_risk = st.tabs([
            "📄 Executive Summary", 
            "📅 Project Timeline", 
            "⚖️ Go/No-Go Analysis",
            "🌐 Translation",
            "🛡️ Risk Management"
        ])

        with tab_summary:
            st.subheader("Auto-Generated Summary")
            summary = generate_summary(extracted_text)
            # Use markdown with class-based styling
            st.markdown(f"<div class='summary-box'>{summary}</div>", unsafe_allow_html=True)
            
            col_a, col_b = st.columns(2)
            with col_a:
                st.subheader("📌 Key Topics Detected")
                if keywords:
                    # Create styled HTML table for key topics
                    table_html = """
                    <table class="styled-table">
                        <thead>
                            <tr>
                                <th>Key Topic</th>
                            </tr>
                        </thead>
                        <tbody>
                    """
                    for kw in keywords:
                        table_html += f"<tr><td>{kw}</td></tr>"
                    table_html += "</tbody></table>"
                    st.markdown(table_html, unsafe_allow_html=True)
                else:
                    st.write("No specific project keywords detected.")
            
            with col_b:
                st.subheader("📅 Dates Mentioned")
                if dates_found:
                    # Create styled HTML table for dates
                    table_html = """
                    <table class="styled-table-dates">
                        <thead>
                            <tr>
                                <th>Date</th>
                            </tr>
                        </thead>
                        <tbody>
                    """
                    for d in dates_found[:10]:
                        table_html += f"<tr><td>{d}</td></tr>"
                    table_html += "</tbody></table>"
                    st.markdown(table_html, unsafe_allow_html=True)
                else:
                    st.write("No specific dates found in document.")

        with tab_timeline:
            st.subheader("Auto-Generated Project Timeline")
            timeline_df = generate_timeline(extracted_text)
            
            fig = px.timeline(
                timeline_df, 
                x_start="Start", 
                x_end="Finish", 
                y="Task", 
                color="Resource", 
                color_discrete_sequence=["#4a5568", "#718096", "#a0aec0", "#cbd5e0"]
            )
            fig.update_yaxes(autorange="reversed")
            fig.update_layout(
                height=400,
                paper_bgcolor='#f8f9fa',
                plot_bgcolor='#ffffff',
                font=dict(size=14),
                margin=dict(l=10, r=10, t=10, b=10)
            )
            st.plotly_chart(fig, use_container_width=True)
            
            st.subheader("Timeline Details")
            # Style the timeline details table
            timeline_html = """
            <table class="timeline-table">
                <thead>
                    <tr>
            """
            for col in timeline_df.columns:
                timeline_html += f"<th>{col}</th>"
            timeline_html += "</tr></thead><tbody>"
            
            for _, row in timeline_df.iterrows():
                timeline_html += "<tr>"
                for col in timeline_df.columns:
                    timeline_html += f"<td>{row[col]}</td>"
                timeline_html += "</tr>"
            
            timeline_html += "</tbody></table>"
            st.markdown(timeline_html, unsafe_allow_html=True)

        with tab_decision:
            c1, c2 = st.columns([2, 1])
            with c1:
                st.subheader("Evaluation Criteria")
                for criterion, score in criteria.items():
                    col_name, col_score, col_bar = st.columns([2, 1, 2])
                    col_name.write(f"**{criterion}:**")
                    col_score.write(f"{score}/10")
                    col_bar.progress(score / 10)
            
            with c2:
                st.subheader("Verdict")
                if "GO ✅" in verdict:
                    st.success(f"# {verdict}")
                elif "NO-GO" in verdict:
                    st.error(f"# {verdict}")
                else:
                    st.warning(f"# {verdict}")
                st.caption(f"Confidence: {confidence}%")

        with tab_translation:
            st.subheader(f"Translate to {target_lang}")
            
            # Get original file type
            file_type = uploaded_file.name.split('.')[-1].lower()
            
            # Translation mode selection
            translation_mode = st.radio(
                "Translation Mode",
                ["Preserve Original Formatting (Recommended)", "Create New Document", "Plain Text Only"],
                horizontal=True,
                help="'Preserve Formatting' keeps the original document structure and styling"
            )
            
            if st.button("🌐 Translate Document", type="primary"):
                original_name = uploaded_file.name.rsplit('.', 1)[0]
                
                if translation_mode == "Preserve Original Formatting (Recommended)":
                    # Translate in-place preserving formatting
                    if file_type == 'docx' and DOCX_SUPPORT:
                        with st.spinner(f'Translating DOCX to {target_lang} (preserving formatting)...'):
                            uploaded_file.seek(0)  # Reset file pointer
                            doc_buffer, error = translate_docx_inplace(uploaded_file, target_lang)
                        
                        if error:
                            st.error(error)
                        else:
                            st.success(f"✅ Document translated with original formatting preserved!")
                            st.download_button(
                                label="📥 Download Translated DOCX (Original Format)",
                                data=doc_buffer,
                                file_name=f"{original_name}_{target_lang.lower().replace(' ', '_')}.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                            )
                    
                    elif file_type in ['pptx', 'ppt']:
                        with st.spinner(f'Translating PPTX to {target_lang} (preserving formatting)...'):
                            uploaded_file.seek(0)  # Reset file pointer
                            pptx_buffer, error = translate_pptx_inplace(uploaded_file, target_lang)
                        
                        if error:
                            st.error(error)
                        else:
                            st.success(f"✅ Presentation translated with original formatting preserved!")
                            st.download_button(
                                label="📥 Download Translated PPTX (Original Format)",
                                data=pptx_buffer,
                                file_name=f"{original_name}_{target_lang.lower().replace(' ', '_')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                    
                    elif file_type == 'pdf':
                        st.warning("⚠️ PDF in-place translation not supported. Using 'Create New Document' mode instead.")
                        # Fall back to creating new document
                        with st.spinner(f'Translating to {target_lang}...'):
                            translated_text = translate_text(extracted_text, target_lang)
                        
                        if DOCX_SUPPORT:
                            doc_buffer = create_translated_docx(translated_text, target_lang, uploaded_file.name)
                            st.download_button(
                                label="📥 Download as Translated DOCX",
                                data=doc_buffer,
                                file_name=f"{original_name}_{target_lang.lower().replace(' ', '_')}.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                            )
                        else:
                            st.download_button(
                                label="📥 Download as Text",
                                data=translated_text,
                                file_name=f"{original_name}_{target_lang.lower()}.txt",
                                mime="text/plain"
                            )
                    else:
                        st.error(f"Unsupported file type for in-place translation: {file_type}")
                
                elif translation_mode == "Create New Document":
                    with st.spinner(f'Translating to {target_lang}...'):
                        translated_text = translate_text(extracted_text, target_lang)
                    
                    st.text_area(
                        f"Translation Preview ({target_lang}):", 
                        translated_text[:2000] + ("..." if len(translated_text) > 2000 else ""), 
                        height=150
                    )
                    
                    # Offer both DOCX and PPTX options
                    col1, col2 = st.columns(2)
                    with col1:
                        if DOCX_SUPPORT:
                            doc_buffer = create_translated_docx(translated_text, target_lang, uploaded_file.name)
                            st.download_button(
                                label="📥 Download as DOCX",
                                data=doc_buffer,
                                file_name=f"{original_name}_{target_lang.lower().replace(' ', '_')}.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                            )
                    with col2:
                        pptx_buffer = create_translated_pptx(translated_text, target_lang, uploaded_file.name)
                        st.download_button(
                            label="📥 Download as PPTX",
                            data=pptx_buffer,
                            file_name=f"{original_name}_{target_lang.lower().replace(' ', '_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    
                    st.success(f"✅ Translation complete! {len(translated_text.split())} words translated.")
                
                else:  # Plain Text Only
                    with st.spinner(f'Translating to {target_lang}...'):
                        translated_text = translate_text(extracted_text, target_lang)
                    
                    st.text_area(
                        f"Translation ({target_lang}):", 
                        translated_text, 
                        height=300
                    )
                    
                    st.download_button(
                        label="📥 Download Translation (TXT)",
                        data=translated_text,
                        file_name=f"{original_name}_{target_lang.lower().replace(' ', '_')}.txt",
                        mime="text/plain"
                    )
                    st.success(f"✅ Translation complete! {len(translated_text.split())} words translated.")
            else:
                st.info("Click the button above to translate the document content.")

        with tab_risk:
            st.subheader("🛡️ Risk Management Dashboard")
            
            # Risk Overview Section
            col_risk1, col_risk2, col_risk3 = st.columns(3)
            
            with col_risk1:
                # Risk Score Gauge
                risk_color = "#43A047" if risk_level == "Low" else "#FFA726" if risk_level == "Medium" else "#EF5350"
                st.metric(
                    label="Overall Risk Score",
                    value=f"{risk_score}/10",
                    delta=f"{risk_level} Risk"
                )
            
            with col_risk2:
                st.metric(
                    label="Project Duration",
                    value=duration,
                    delta="Estimated"
                )
            
            with col_risk3:
                st.metric(
                    label="Confidence Level",
                    value=f"{confidence}%",
                    delta="Analysis Confidence"
                )
            
            st.markdown("---")
            
            # Risk Categories Analysis
            st.subheader("📊 Risk Categories Breakdown")
            
            # Extract risk-related keywords for categorization
            text_lower = extracted_text.lower()
            
            risk_categories = {
                "Technical Risk": {
                    "indicators": ["technical", "system", "software", "hardware", "integration", "compatibility", "bug", "error", "failure"],
                    "score": 0
                },
                "Financial Risk": {
                    "indicators": ["budget", "cost", "expense", "funding", "investment", "financial", "money", "price", "revenue"],
                    "score": 0
                },
                "Schedule Risk": {
                    "indicators": ["deadline", "delay", "timeline", "schedule", "milestone", "late", "overdue", "time"],
                    "score": 0
                },
                "Resource Risk": {
                    "indicators": ["resource", "staff", "team", "personnel", "capacity", "availability", "shortage", "skill"],
                    "score": 0
                },
                "Compliance Risk": {
                    "indicators": ["compliance", "regulatory", "legal", "policy", "standard", "requirement", "audit", "governance"],
                    "score": 0
                },
                "Operational Risk": {
                    "indicators": ["operational", "process", "workflow", "efficiency", "performance", "quality", "maintenance"],
                    "score": 0
                }
            }
            
            # Calculate scores for each category
            for category, data in risk_categories.items():
                count = sum(1 for word in data["indicators"] if word in text_lower)
                data["score"] = min(10, count * 2)  # Scale to 0-10
            
            # Display risk categories
            col_cat1, col_cat2 = st.columns(2)
            
            categories_list = list(risk_categories.items())
            
            with col_cat1:
                for category, data in categories_list[:3]:
                    score = data["score"]
                    level = "🟢 Low" if score <= 3 else "🟡 Medium" if score <= 6 else "🔴 High"
                    st.write(f"**{category}:** {level}")
                    st.progress(score / 10)
            
            with col_cat2:
                for category, data in categories_list[3:]:
                    score = data["score"]
                    level = "🟢 Low" if score <= 3 else "🟡 Medium" if score <= 6 else "🔴 High"
                    st.write(f"**{category}:** {level}")
                    st.progress(score / 10)
            
            st.markdown("---")
            
            # Risk Indicators Found
            st.subheader("⚠️ Risk Indicators Detected")
            
            high_risk_words = ['urgent', 'critical', 'risk', 'delay', 'issue', 'problem', 'challenge', 'concern', 'failure', 'crisis', 'threat', 'danger', 'warning']
            medium_risk_words = ['consider', 'review', 'assess', 'evaluate', 'potential', 'uncertain', 'unclear', 'possible', 'might', 'could']
            positive_words = ['complete', 'success', 'achieved', 'approved', 'confirmed', 'stable', 'secure', 'resolved', 'mitigated', 'controlled']
            
            found_high = [word for word in high_risk_words if word in text_lower]
            found_medium = [word for word in medium_risk_words if word in text_lower]
            found_positive = [word for word in positive_words if word in text_lower]
            
            col_ind1, col_ind2, col_ind3 = st.columns(3)
            
            with col_ind1:
                st.markdown("**🔴 High Risk Indicators**")
                if found_high:
                    for word in found_high[:5]:
                        st.write(f"• {word.capitalize()}")
                else:
                    st.write("None detected ✓")
            
            with col_ind2:
                st.markdown("**🟡 Medium Risk Indicators**")
                if found_medium:
                    for word in found_medium[:5]:
                        st.write(f"• {word.capitalize()}")
                else:
                    st.write("None detected ✓")
            
            with col_ind3:
                st.markdown("**🟢 Positive Indicators**")
                if found_positive:
                    for word in found_positive[:5]:
                        st.write(f"• {word.capitalize()}")
                else:
                    st.write("None detected")
            
            st.markdown("---")
            
            # Risk Mitigation Recommendations
            st.subheader("💡 Risk Mitigation Recommendations")
            
            recommendations = []
            
            if risk_score < 4:
                recommendations.append("⚠️ **High Priority:** Conduct immediate risk assessment meeting with stakeholders")
                recommendations.append("📋 Develop detailed contingency plans for identified risks")
                recommendations.append("👥 Consider allocating additional resources to risk mitigation")
            
            if any(cat["score"] > 6 for cat in risk_categories.values()):
                high_risk_cats = [cat for cat, data in risk_categories.items() if data["score"] > 6]
                for cat in high_risk_cats:
                    recommendations.append(f"🎯 Focus on reducing **{cat}** through targeted interventions")
            
            if "delay" in text_lower or "deadline" in text_lower:
                recommendations.append("⏰ Implement schedule monitoring and early warning systems")
            
            if "budget" in text_lower or "cost" in text_lower:
                recommendations.append("💰 Establish budget tracking and cost control measures")
            
            if len(recommendations) == 0:
                recommendations.append("✅ Risk levels appear manageable - maintain regular monitoring")
                recommendations.append("📊 Continue periodic risk assessments throughout project lifecycle")
                recommendations.append("📝 Document lessons learned for future risk management")
            
            for rec in recommendations:
                st.markdown(rec)
            
            st.markdown("---")
            
            # Risk Summary Table
            st.subheader("📋 Risk Summary Report")
            
            risk_summary_html = f"""
            <table class="styled-table">
                <thead>
                    <tr>
                        <th>Metric</th>
                        <th>Value</th>
                        <th>Status</th>
                    </tr>
                </thead>
                <tbody>
                    <tr>
                        <td>Overall Risk Score</td>
                        <td>{risk_score}/10</td>
                        <td>{'🟢' if risk_score >= 7 else '🟡' if risk_score >= 4 else '🔴'} {risk_level}</td>
                    </tr>
                    <tr>
                        <td>High Risk Indicators</td>
                        <td>{len(found_high)}</td>
                        <td>{'🟢' if len(found_high) <= 1 else '🟡' if len(found_high) <= 3 else '🔴'}</td>
                    </tr>
                    <tr>
                        <td>Positive Indicators</td>
                        <td>{len(found_positive)}</td>
                        <td>{'🟢' if len(found_positive) >= 3 else '🟡' if len(found_positive) >= 1 else '🔴'}</td>
                    </tr>
                    <tr>
                        <td>Go/No-Go Decision</td>
                        <td>{verdict}</td>
                        <td>{confidence}% Confidence</td>
                    </tr>
                    <tr>
                        <td>Estimated Duration</td>
                        <td>{duration}</td>
                        <td>📅</td>
                    </tr>
                </tbody>
            </table>
            """
            st.markdown(risk_summary_html, unsafe_allow_html=True)